from concurrent.futures import ThreadPoolExecutor
from flask import Flask, request, send_file
import openpyxl
from io import BytesIO
import base64
from pptx import Presentation
from docx import Document
from googletrans import Translator as trans
# import PyMuPDF
# import fitz  # PyMuPDF
# from PyPDF2 import PdfReader, PdfWriter
import os
import pandas as pd
# import PyPDF2
# from tika import parser
# from pdf2docx import Converter
# import docx2pdf
# from docx2pdf import convert as docx2pdf_convert
from fpdf import FPDF

import pdfplumber
from functools import partial
from nltk.tokenize import sent_tokenize
# import nltk
from deep_translator import GoogleTranslator
import nltk
nltk.download('punkt')

app = Flask(__name__)
 
translator = trans()

def not_within_bboxes(obj, bboxes):
    """Check if the object is in any of the table's bbox."""
    def obj_in_bbox(_bbox):
        v_mid = (obj["top"] + obj["bottom"]) / 2
        h_mid = (obj["x0"] + obj["x1"]) / 2
        x0, top, x1, bottom = _bbox
        return (h_mid >= x0) and (h_mid < x1) and (v_mid >= top) and (v_mid < bottom)
    return not any(obj_in_bbox(__bbox) for __bbox in bboxes)

def extract(page):
    """Extract PDF text.

    Filter out tables and delete in-paragraph line-breaks.
    """
    # Filter-out tables
    if page.find_tables() != []:
        # Get the bounding boxes of the tables on the page.
        bboxes = [table.bbox for table in page.find_tables()]
        
        bbox_not_within_bboxes = partial(not_within_bboxes, bboxes=bboxes) 

        # Filter-out tables from page
        page = page.filter(bbox_not_within_bboxes)

    # Extract Text
    extracted = page.extract_text()

    # Delete in-paragraph line breaks
    extracted = extracted.replace(".\n", "/m" # keep paragraph breaks
                        ).replace(". \n", "/m" # keep paragraph breaks
                        ).replace("\n", "" # delete in-paragraph breaks (i.e. all remaining \n)
                        ).replace("/m", ".\n\n") # restore paragraph breaks
    
    return extracted

def translate_extracted(Extracted,langaugeFrom,langaugeTo):
    """Wrapper for Google Translate with upload workaround.
    
    Collects chuncks of senteces below limit to translate.
    """
    # Set-up and wrap translation client
    translate = GoogleTranslator(source=langaugeFrom, target=langaugeTo).translate

    # Split input text into a list of sentences
    sentences = sent_tokenize(Extracted)

    # Initialize containers
    translated_text = ''
    source_text_chunk = ''

    # collect chuncks of sentences below limit and translate them individually
    for sentence in sentences:
        # if chunck together with current sentence is below limit, add the sentence
        if ((len(sentence.encode('utf-8')) + len(source_text_chunk.encode('utf-8')) < 5000)):
            source_text_chunk += ' ' + sentence
        
        # else translate chunck and start new one with current sentence
        else:
            translated_text += ' ' + translate(source_text_chunk)

            # if current sentence smaller than 5000 chars, start new chunck
            if (len(sentence.encode('utf-8')) < 5000):
                source_text_chunk = sentence

            # else, replace sentence with notification message
            else:    
                message = "<<Omitted Word longer than 5000bytes>>"
                translated_text += ' ' + translate(message)

                # Re-set text container to empty
                source_text_chunk = ''

    # Translate the final chunk of input text, if there is any valid text left to translate
    if translate(source_text_chunk) != None:
        translated_text += ' ' + translate(source_text_chunk)
    
    return translated_text

def annotate_pdf(pdf_path, changes):
    # Modified function to annotate based on the changes provided
    pdf_document = fitz.open(pdf_path)
    for page_num in range(len(pdf_document)):
        page = pdf_document[page_num]
        print("splited data::",changes.split(" "))
        for text in changes.split(" "):
            quads = page.search_for(text, quads=True)
            if quads:
                page.add_highlight_annot(quads)
    annotated_pdf_path = pdf_path.replace('.pdf', '_annotated.pdf')
    pdf_document.save(annotated_pdf_path)
    pdf_document.close()
    print(f"Annotated PDF saved as: {annotated_pdf_path}")

# def translate_text(text, target_language='en'):
#     """Translate text using Google Translate."""
#     print(target_language)
#     translated_text = translator.translate(text,src="ja", dest="en")
#     return translated_text.text

def translate_text(text, target_language='en'):
    """Translate text using Google Translate."""
    try:
        print(target_language)
        translated_text = translator.translate(text, src="ja", dest="en")
        translated_text_str = translated_text.text.encode('utf-8').decode('utf-8')
        print(translated_text_str)
        return translated_text_str
    except Exception as error:
        print(f"Translation error: {str(error)}")
        return None


def save_uploaded_file(uploaded_file):
    if uploaded_file is not None:
        path = os.getcwd() + "/client/src/pdf/"
        print("working directory::",path)
        file_path = os.path.join(path, "test.pdf")  # Save file to a temporary directory
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        return file_path
    else:
        return None
    
def process_pdfs(pdf1_path):
  # Extract content from PDFs
  text1, tables1 = extract_content_from_pdf(pdf1_path)

  text1 = text1.replace('\n', '')
  translated_text = translate_text(text1,"ja")
#   print(translated_text)
  return translated_text
#   differ = difflib.Differ()
#   diff_me = list(differ.compare(text1, text2))
    

def extract_content_from_pdf(pdf_path):
    """
    Extract text and tables from a PDF file.
    """
    try:
        text = ''
        tables = []
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
                page_tables = page.extract_tables()
                for table in page_tables:
                    tables.append(pd.DataFrame(table[1:], columns=table[0]))
        return text, tables
    except Exception as e:
        return None, None

def translate_pdf(pdf_content, language_to, language_from):
    try:
        translated_pdf_content = BytesIO()

        with fitz.open(stream=pdf_content, filetype="pdf") as pdf_document:
            pdf_writer = fitz.open()

            for page_number in range(pdf_document.page_count):
                page = pdf_document[page_number]
                text = page.get_text()

                translated_text = translator.translate(text, dest=language_to, src=language_from).text
                print(translated_text)
                translated_page = pdf_writer.new_page(width=page.rect.width, height=page.rect.height)
                
                # Insert the translated text onto the new page
                translated_page.insert_text((10, 10), translated_text)

            pdf_writer.save(translated_pdf_content, garbage=3, deflate=True)
        
        translated_pdf_content.seek(0)
        return translated_pdf_content.getvalue()

    except Exception as error:
        print(f"Translation error: {str(error)}")
        return None

def translate_cell(cell, language_to, language_from):
    try:
        if cell.value is not None:
            
            if hasattr(cell.value, 'iter') and not isinstance(cell.value, str):
                
                translated_values = [translator.translate(str(value), dest=language_to, src=language_from).text for value in cell.value]
                cell.value = translated_values
            elif isinstance(cell.value, (str, int, float)):
                
                translated_text = translator.translate(str(cell.value), dest=language_to, src=language_from).text
                cell.value = translated_text
    except TypeError:
        # Handle 'NoneType' error explicitly and skip translation for None values
        pass
    except Exception as error:
        print(f"Translation error: {str(error)}")


def translate_text_block(text_frame, language_to, language_from):
    try:
        if text_frame is not None:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        translated_text = translator.translate(run.text, dest=language_to, src=language_from).text
                        run.text = translated_text
    except Exception as error:
        print(f"Translation error block: {str(error)}")

def translate_docx_paragraph(paragraph, language_to, language_from):
    try:
        for run in paragraph.runs:
            try:
                if run.text:
                    translated_text = translator.translate(run.text, dest=language_to, src=language_from).text
                    print(run.text ,"-" , translated_text)
                    run.text = translated_text
            except:
                continue
    except Exception as error:
        print(f"Translation error paragrph: {str(error)}")

def translate_docx_table(table, language_to, language_from):
    try:
        for row in table.rows:
            for cell in row.cells:
                if cell.paragraphs[0] != None:
                    translate_docx_paragraph(cell.paragraphs[0], language_to, language_from)
    except Exception as error:
        print(f"Translation error Table: {str(error)}")
        
        
def convert_docx_to_pdf(docx_path, pdf_path):
    document = Document(docx_path)

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    for paragraph in document.paragraphs:
        try:
            # Initialize font size and font family
            font_size = 12
            font_family = "Arial"
            para = ""
            for run in paragraph.runs:

                # Try applying the font for the current run
                try:
                    pdf.set_font(font_family, size=font_size)
                    para = para + " "+ run.text  # Append text from the current run
                except Exception as font_error:
                    print(f"Font error for run '{run.text}': {font_error}")

            # Use the original font size and font family in the PDF
            pdf.set_font(font_family, size=font_size)

            # Adjust the width of the cell based on the available width of the page
            cell_width = pdf.w - 2 * pdf.l_margin

            # Set x-position to left margin before writing each cell
            pdf.set_x(pdf.l_margin)

            # Calculate width dynamically based on content
            cell_width = pdf.get_string_width(para) + 6  # 6 is added for padding

            # Multi-cell with adjusted width
            pdf.multi_cell(cell_width, 10, para)
        except Exception as e:
            print(f"Error processing paragraph: {e}")
            continue

    pdf.output(pdf_path)

@app.route('/convert', methods=['POST'])
def convert():
    try:
        file = request.files['file']
        language_to = request.form.get('languageTo', 'en')
        language_from = request.form.get('languageFrom', 'en')
        file_extension = request.form.get("fileExtension", 'en')

        if file:
            if language_to == "":
                language_to = "en"

            file_extension = file_extension.split(".")[-1]
            if file_extension == "xlsx":
                file_content_base64 = request.form.get('fileContent')
                decoded_content = base64.b64decode(file_content_base64)
                workbook = openpyxl.load_workbook(BytesIO(decoded_content))

                with ThreadPoolExecutor() as executor:
                    futures = []
                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        for row in sheet.iter_rows():
                            for cell in row:
                                future = executor.submit(translate_cell, cell, language_to, language_from)
                                futures.append(future)

                    for future in futures:
                        future.result()

                translated_content = BytesIO()
                workbook.save(translated_content)
                translated_content.seek(0)

                return send_file(
                    translated_content,
                    mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                    as_attachment=True,
                    download_name=file.name
                )

            elif file_extension == "pptx":
                file_content_base64 = request.form.get('fileContent')
                decoded_content = base64.b64decode(file_content_base64)
                presentation = Presentation(BytesIO(decoded_content))

                with ThreadPoolExecutor() as executor:
                    futures = []
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                future = executor.submit(translate_text_block, shape.text_frame, language_to, language_from)
                                futures.append(future)

                    for future in futures:
                        future.result()

                translated_content = BytesIO()
                presentation.save(translated_content)
                translated_content.seek(0)

                return send_file(
                    translated_content,
                    mimetype='application/octet-stream',
                    as_attachment=True,
                    download_name='translated_presentation.pptx'
                )

            elif file_extension == "docx":
                
                file_content_base64 = request.form.get('fileContent')
                decoded_content = base64.b64decode(file_content_base64)
                document = Document(BytesIO(decoded_content))

                with ThreadPoolExecutor() as executor:
                    futures = []
                    for paragraph in document.paragraphs:
                        future = executor.submit(translate_docx_paragraph, paragraph, language_to, language_from)
                        futures.append(future)

                    for table in document.tables:
                        future = executor.submit(translate_docx_table, table, language_to, language_from)
                        futures.append(future)

                    for future in futures:
                        future.result()

                translated_content = BytesIO()
                document.save(translated_content)
                translated_content.seek(0)

                return send_file(
                    translated_content,
                    mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    as_attachment=True,
                    download_name=file.name
                )
            elif file_extension == "pdf":
                print("language_from:", language_from,", language_to::",language_to)
                file_path = save_uploaded_file(file)
                saved_pdf_path = os.getcwd() + "/client/src/pdf/output.pdf"
                with pdfplumber.open(file_path) as pdf:
                    # Initialize FPDF file to write on
                    fpdf = FPDF()
                    fpdf.set_font("Helvetica", size = 10)
                    
                    # Treat each page individually
                    for page in pdf.pages:
                        # Extract Page
                        extracted = extract(page)

                        # Translate Page
                        if extracted != "":
                            # Translate paragraphs individually to keep breaks
                            paragraphs = extracted.split("\n\n")
                            translated = "\n\n".join(
                                [translate_extracted(paragraph,language_from,language_to) for paragraph in paragraphs]
                                )
                        else:
                            translated = extracted

                        # Write Page
                        fpdf.add_page()
                        fpdf.multi_cell(w=0,
                                        h=5,
                                        txt= translated.encode("latin-1",
                                                            errors = "replace"
                                                    ).decode("latin-1")
                                        )
                    
                    # Save all FPDF pages
                    fpdf.output(saved_pdf_path)
                return send_file(
                    os.getcwd() + "/client/src/pdf/output.pdf",
                    mimetype='application/pdf',
                    as_attachment=True,
                    download_name='output.pdf'
                )                          
            
    except Exception as e:
        print("error:",e)

    
    return {'message': 'File conversion successful'}, 200

if __name__ == "__main__":
    app.run(debug=True)